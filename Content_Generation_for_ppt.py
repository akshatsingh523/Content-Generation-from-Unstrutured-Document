import logging
import os
import time
import asyncio
from datetime import datetime
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import io
from dotenv import load_dotenv
import openai
from fastapi import FastAPI, UploadFile, File
from typing import List
from fastapi.responses import JSONResponse
import google.generativeai as genai

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Load environment variables from .env file
load_dotenv()
app = FastAPI()

# Configure OpenAI API Key
api_key = os.getenv("OPENAI_API_KEY")
openai.api_key = api_key

# Configure Google API Key for Gemini Pro Vision
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))

safety_settings = [
    {
        "category": "HARM_CATEGORY_DANGEROUS",
        "threshold": "BLOCK_NONE",
    },
    {
        "category": "HARM_CATEGORY_HARASSMENT",
        "threshold": "BLOCK_NONE",
    },
    {
        "category": "HARM_CATEGORY_HATE_SPEECH",
        "threshold": "BLOCK_NONE",
    },
    {
        "category": "HARM_CATEGORY_SEXUALLY_EXPLICIT",
        "threshold": "BLOCK_NONE",
    },
    {
        "category": "HARM_CATEGORY_DANGEROUS_CONTENT",
        "threshold": "BLOCK_NONE",
    },
]

# Load Gemini pro vision model
model = genai.GenerativeModel(model_name='gemini-pro-vision', safety_settings=safety_settings)

# Constants for cost calculation (example values, adjust based on actual pricing)
OPENAI_COST_PER_TOKEN = 0.00003  # Cost per token processed (example value)

# Constants for Google Gemini API
GOOGLE_GEMINI_COST_INPUT_SHORT = 0.35 / 1_000_000
GOOGLE_GEMINI_COST_INPUT_LONG = 0.70 / 1_000_000
GOOGLE_GEMINI_COST_OUTPUT_SHORT = 1.05 / 1_000_000
GOOGLE_GEMINI_COST_OUTPUT_LONG = 2.10 / 1_000_000

def calculate_openai_cost(num_tokens):
    return num_tokens * OPENAI_COST_PER_TOKEN

def calculate_google_gemini_cost(input_tokens, output_tokens):
    if input_tokens <= 128_000:
        input_cost = input_tokens * GOOGLE_GEMINI_COST_INPUT_SHORT
        output_cost = output_tokens * GOOGLE_GEMINI_COST_OUTPUT_SHORT
    else:
        input_cost = input_tokens * GOOGLE_GEMINI_COST_INPUT_LONG
        output_cost = output_tokens * GOOGLE_GEMINI_COST_OUTPUT_LONG
    return input_cost + output_cost

def get_token_count(text):
    # Simple token count based on whitespace
    return len(text.split())

def get_summary(prompt):
    try:
        response = openai.ChatCompletion.create(
            model="gpt-4-turbo",
            messages=[
                {"role": "system", "content": "You are a helpful assistant."},
                {"role": "user", "content": prompt}
            ],
            max_tokens=300
        )
        return response.choices[0].message['content'].strip()
    except Exception as e:
        logger.error(f"Error in get_summary: {e}")
        return ""

def extract_ppt_content(ppt_file):
    if not os.path.exists(ppt_file):
        raise FileNotFoundError(f"Package not found at '{ppt_file}'")

    presentation = Presentation(ppt_file)
    content = []
    
    image_counter = 1
    
    for slide_number, slide in enumerate(presentation.slides):
        slide_content = {
            "slide_number": slide_number + 1,
            "texts": [],
            "tables": [],
            "images": []
        }
        
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for shp in shape.shapes:
                    image_counter = process_shape(shp, slide_content, image_counter)
            else:
                image_counter = process_shape(shape, slide_content, image_counter)
                
        content.append(slide_content)
    
    return content

def process_shape(shape, slide_content, image_counter):
    try:
        if shape.has_text_frame:
            text = shape.text_frame.text
            slide_content["texts"].append(text)

        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table_content = []
            table = shape.table
            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    row_data.append(cell.text)
                table_content.append(row_data)
            slide_content["tables"].append(table_content)

        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image = shape.image
            image_bytes = io.BytesIO(image.blob)
            img = Image.open(image_bytes)

            logger.info(f"Processing image {image_counter} on slide {slide_content['slide_number']}")

            # Use Gemini Pro Vision to analyze the image directly as a PIL Image
            response = model.generate_content([img])

            if response and response.candidates:
                candidate = response.candidates[0]

                # Extract text from the parts attribute
                description = ""
                if hasattr(candidate, 'content') and hasattr(candidate.content, 'parts'):
                    for part in candidate.content.parts:
                        if hasattr(part, 'text'):
                            description = part.text
                            break

                if description:
                    image_info = {
                        "description": description
                    }
                    slide_content["images"].append(image_info)
                    logger.info(f"Added image description: {description}")
                else:
                    logger.warning(f"No valid description found for image on slide {slide_content['slide_number']}")
            else:
                logger.warning(f"Image description was blocked or no valid response for image on slide {slide_content['slide_number']}")

            # Increment the image counter
            image_counter += 1

        return image_counter
    except Exception as e:
        logger.error(f"Error in process_shape: {e}")
        raise

def preprocess_ppt_content(content):
    try:
        full_extracted_content = []
        combined_text = ""
        
        for slide in content:
            slide_info = {
                "slide_number": slide["slide_number"],
                "texts": [],
                "tables": [],
                "image_descriptions": []
            }

            # Process texts
            for text in slide["texts"]:
                slide_info["texts"].append(text)
                combined_text += text + "\n"
            
            # Process tables
            for table in slide["tables"]:
                table_text = "\n".join(["\t".join(row) for row in table])
                slide_info["tables"].append(table_text)
                combined_text += table_text + "\n"
            
            # Process images
            for image_info in slide["images"]:
                image_description = f"[Image Description: {image_info['description']}]"
                slide_info["image_descriptions"].append(image_description)
                combined_text += image_description + "\n"
            
            full_extracted_content.append(slide_info)
        
        return full_extracted_content, combined_text
    except Exception as e:
        logger.error(f"Error in preprocess_ppt_content: {e}")
        raise

def generate_openai_responses(input_text, prompt_type):
    prompts = {
        "elevator_pitch": f"Summarize the following content into a concise and compelling elevator pitch. Focus on capturing the essence and key points in a conversational tone, suitable for quick presentations and pitches. Ensure the summary is engaging and highlights the most important aspects. Here is the content: {input_text}",
        "short_description": f"Provide a brief description of the following content. The summary should be short and clear, giving an overview of the main points and themes. Aim to provide a quick understanding of the key elements without going into too much detail. Here is the content: {input_text}",
        "long_description": f"Create a detailed description of the following content. Cover all critical aspects and details, ensuring that the summary is thorough and comprehensive. Include important points, themes, and any relevant specifics that provide a deep understanding of the content. Here is the content: {input_text}",
        "voice_over": f"Develop a voice-over script for the following content, suitable for use in an explainer video. Emphasize key points and messages, ensuring a smooth and conversational flow. Avoid using bullet points; instead, create a narrative that guides the viewer through the content in an engaging and informative manner. Here is the content: {input_text}"
    }
    
    try:
        response = openai.ChatCompletion.create(
            model="gpt-4-turbo",
            messages=[
                {"role": "system", "content": "You are a helpful assistant."},
                {"role": "user", "content": prompts[prompt_type]}
            ],
            max_tokens=300
        )
        return response.choices[0].message['content'].strip()
    except Exception as e:
        logger.error(f"Error in generate_openai_responses: {e}")
        return ""

def generate_slide_summary(content, prompt_type):
    try:
        slide_summaries = []
        for slide in content:
            print(f"Processing slide {slide['slide_number']} at {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
            slide_text = "\n".join(slide["texts"] + slide["tables"] + slide["image_descriptions"])
            response = generate_openai_responses(slide_text, prompt_type)
            slide_summaries.append({
                "slide_number": slide["slide_number"],
                "summary": response
            })
            print(f"Slide {slide['slide_number']} Summary: {response}")
            print(f"Finished processing slide {slide['slide_number']} at {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
            # Adding sleep to respect rate limits
            time.sleep(1)
        return slide_summaries
    except Exception as e:
        logger.error(f"Error in generate_slide_summary: {e}")
        raise

def generate_overall_summary(combined_text, prompt_type):
    try:
        response = generate_openai_responses(combined_text, prompt_type)
        print(f"Overall Summary: {response}")
        return response
    except Exception as e:
        logger.error(f"Error in generate_overall_summary: {e}")
        raise

# Functions for Elevator Pitch
def generate_slide_elevator_pitch(content):
    return generate_slide_summary(content, "elevator_pitch")

def generate_overall_elevator_pitch(combined_text):
    return generate_overall_summary(combined_text, "elevator_pitch")

# Functions for Short Description
def generate_slide_short_description(content):
    return generate_slide_summary(content, "short_description")

def generate_overall_short_description(combined_text):
    return generate_overall_summary(combined_text, "short_description")

# Functions for Long Description
def generate_slide_long_description(content):
    return generate_slide_summary(content, "long_description")

def generate_overall_long_description(combined_text):
    return generate_overall_summary(combined_text, "long_description")

# Functions for Voice Over
def generate_slide_voice_over(content):
    return generate_slide_summary(content, "voice_over")

def generate_overall_voice_over(combined_text):
    return generate_overall_summary(combined_text, "voice_over")

async def process_upload_ppt(file: UploadFile, timeout: int):
    start_time = time.time()
    logger.info(f"Processing started at {time.strftime('%Y-%m-%d %H:%M:%S', time.localtime(start_time))}")
    try:
        file_location = f"temp_{file.filename}"
        with open(file_location, "wb+") as file_object:
            file_object.write(file.file.read())
        
        logger.info("File uploaded successfully.")
        
        # Use asyncio.wait_for to add a timeout
        content = await asyncio.wait_for(asyncio.to_thread(extract_ppt_content, file_location), timeout)
        full_content, combined_text = await asyncio.wait_for(asyncio.to_thread(preprocess_ppt_content, content), timeout)
        
        os.remove(file_location)

        logger.info(f"Full content extracted: {full_content}")

        # Calculate the costs
        num_images = sum(len(slide.get('images', [])) for slide in full_content)
        logger.info(f"Number of images processed: {num_images}")
        
        num_input_tokens = get_token_count(combined_text)
        num_output_tokens = num_input_tokens  # Assuming output tokens are similar in count to input tokens for simplicity
        google_gemini_cost = calculate_google_gemini_cost(num_input_tokens, num_output_tokens)
        
        num_openai_tokens = get_token_count(combined_text)
        openai_cost = calculate_openai_cost(num_openai_tokens)
        
        total_cost = google_gemini_cost + openai_cost
        
        end_time = time.time()
        logger.info(f"Processing finished at {time.strftime('%Y-%m-%d %H:%M:%S', time.localtime(end_time))}")
        logger.info(f"Total processing time: {end_time - start_time:.2f} seconds")
        
        return {
            "message": "File processed successfully",
            "content": full_content,
            "combined_text": combined_text,
            "google_gemini_cost": google_gemini_cost,
            "openai_cost": openai_cost,
            "total_cost": total_cost
        }
    except asyncio.TimeoutError:
        logger.error("Request timed out.")
        return {"message": "Request timed out"}
    except Exception as e:
        logger.error(f"Error occurred: {str(e)}")
        return {"message": str(e)}

@app.get("/")
async def read_root():
    return {"message": "Hello, FastAPI!"}

@app.post("/upload_ppt/")
async def upload_ppt(file: UploadFile = File(...)):
    return await process_upload_ppt(file, timeout=3600)

@app.post("/elevator_pitch/")
async def elevator_pitch(file: UploadFile = File(...)):
    try:
        response = await upload_ppt(file)
        if response["message"] == "File processed successfully":
            full_content = response["content"]
            combined_text = response["combined_text"]

            slide_elevator_pitches = generate_slide_elevator_pitch(full_content)
            overall_elevator_pitch = generate_overall_elevator_pitch(combined_text)

            return {
                "google_gemini_cost": response["google_gemini_cost"],
                "openai_cost": response["openai_cost"],
                "total_cost": response["total_cost"],
                "slide_elevator_pitches": slide_elevator_pitches,
                "overall_elevator_pitch": overall_elevator_pitch
            }
        else:
            return JSONResponse(status_code=500, content={"message": response["message"]})
    except Exception as e:
        logger.error(f"Error in elevator_pitch: {e}")
        return JSONResponse(status_code=500, content={"message": str(e)})

@app.post("/short_description/")
async def short_description(file: UploadFile = File(...)):
    try:
        response = await upload_ppt(file)
        if response["message"] == "File processed successfully":
            full_content = response["content"]
            combined_text = response["combined_text"]

            slide_short_descriptions = generate_slide_short_description(full_content)
            overall_short_description = generate_overall_short_description(combined_text)

            return {
                "google_gemini_cost": response["google_gemini_cost"],
                "openai_cost": response["openai_cost"],
                "total_cost": response["total_cost"],
                "slide_short_descriptions": slide_short_descriptions,
                "overall_short_description": overall_short_description
            }
        else:
            return JSONResponse(status_code=500, content={"message": response["message"]})
    except Exception as e:
        logger.error(f"Error in short_description: {e}")
        return JSONResponse(status_code=500, content={"message": str(e)})

@app.post("/long_description/")
async def long_description(file: UploadFile = File(...)):
    try:
        response = await upload_ppt(file)
        if response["message"] == "File processed successfully":
            full_content = response["content"]
            combined_text = response["combined_text"]

            slide_long_descriptions = generate_slide_long_description(full_content)
            overall_long_description = generate_overall_long_description(combined_text)

            return {
                "google_gemini_cost": response["google_gemini_cost"],
                "openai_cost": response["openai_cost"],
                "total_cost": response["total_cost"],
                "slide_long_descriptions": slide_long_descriptions,
                "overall_long_description": overall_long_description
            }
        else:
            return JSONResponse(status_code=500, content={"message": response["message"]})
    except Exception as e:
        logger.error(f"Error in long_description: {e}")
        return JSONResponse(status_code=500, content={"message": str(e)})

@app.post("/voice_over/")
async def voice_over(file: UploadFile = File(...)):
    try:
        response = await upload_ppt(file)
        if response["message"] == "File processed successfully":
            full_content = response["content"]
            combined_text = response["combined_text"]

            slide_voice_overs = generate_slide_voice_over(full_content)
            overall_voice_over = generate_overall_voice_over(combined_text)

            return {
                "google_gemini_cost": response["google_gemini_cost"],
                "openai_cost": response["openai_cost"],
                "total_cost": response["total_cost"],
                "slide_voice_overs": slide_voice_overs,
                "overall_voice_over": overall_voice_over
            }
        else:
            return JSONResponse(status_code=500, content={"message": response["message"]})
    except Exception as e:
        logger.error(f"Error in voice_over: {e}")
        return JSONResponse(status_code=500, content={"message": str(e)})

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000, timeout_keep_alive=3600)
